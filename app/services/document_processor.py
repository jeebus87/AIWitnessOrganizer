"""Document processing service for PDFs, images, and Outlook emails"""
import logging
import os
import io
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
import hashlib

from PIL import Image
import extract_msg

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

logger = logging.getLogger(__name__)


@dataclass
class ProcessedAsset:
    """Represents a processed document asset ready for AI analysis"""
    asset_type: str  # "image", "text", "email_body"
    content: bytes  # Image bytes or text encoded as bytes
    media_type: str  # "image/jpeg", "image/png", "text/plain"
    filename: str
    original_filename: str
    parent_filename: Optional[str] = None
    context: str = ""  # Additional context (e.g., "Attachment to email from X")
    page_number: Optional[int] = None  # For multi-page PDFs


@dataclass
class ProcessingResult:
    """Result of processing a document"""
    success: bool
    assets: List[ProcessedAsset] = field(default_factory=list)
    error: Optional[str] = None
    file_hash: Optional[str] = None


class DocumentProcessor:
    """
    Processes various document types into assets ready for AI analysis.
    Supports PDFs, images, and Outlook .msg/.eml files with recursive extraction.
    """

    # Image constraints for AWS Bedrock
    # Multi-image requests (30 pages per chunk) have a 2000px limit per image
    MAX_IMAGE_SIZE_MB = 3.75
    MAX_IMAGE_DIMENSION = 1920  # Under 2000px limit for multi-image requests
    SUPPORTED_IMAGE_FORMATS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tif", ".tiff"}
    SUPPORTED_DOC_FORMATS = {".pdf", ".msg", ".eml", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".rtf", ".txt", ".html", ".htm", ".csv"}

    def __init__(self, temp_dir: Optional[str] = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()

    def get_file_hash(self, content: bytes) -> str:
        """Calculate SHA-256 hash of file content for caching"""
        return hashlib.sha256(content).hexdigest()

    def detect_file_type(self, filename: str, content: bytes) -> str:
        """Detect file type from extension and content"""
        ext = Path(filename).suffix.lower()

        # Check magic bytes for common formats
        if content[:4] == b"%PDF":
            return "pdf"
        elif content[:8] == b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1":  # OLE compound (could be .msg, .doc, .xls, .ppt)
            # Need to check extension for legacy Office formats
            if ext == ".msg":
                return "msg"
            elif ext == ".doc":
                return "doc"
            elif ext == ".xls":
                return "xls"
            elif ext == ".ppt":
                return "ppt"
            return "msg"  # Default to msg for OLE
        elif content[:4] == b"PK\x03\x04":  # ZIP-based (Office Open XML: .docx, .xlsx, .pptx)
            if ext == ".docx":
                return "docx"
            elif ext in (".xlsx", ".xls"):
                return "xlsx"
            elif ext == ".pptx":
                return "pptx"
            # Could also be a regular .zip, fall back to extension
            return ext.lstrip(".") if ext else "zip"
        elif content[:4] == b"\x89PNG":
            return "png"
        elif content[:2] == b"\xFF\xD8":
            return "jpg"
        elif content[:4] in (b"II*\x00", b"MM\x00*"):  # TIFF (little/big endian)
            return "tiff"
        elif content[:5] == b"{\\rtf":
            return "rtf"
        elif content[:14].lower().startswith(b"<!doctype html") or content[:5].lower() == b"<html":
            return "html"
        elif ext in self.SUPPORTED_IMAGE_FORMATS:
            return ext.lstrip(".")

        return ext.lstrip(".") or "unknown"

    async def process(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> ProcessingResult:
        """
        Process a document and return assets ready for AI analysis.

        Args:
            content: Raw file bytes
            filename: Original filename
            context: Additional context for nested files

        Returns:
            ProcessingResult with list of ProcessedAssets
        """
        file_hash = self.get_file_hash(content)
        file_type = self.detect_file_type(filename, content)
        
        logger.info(f"Processing document: {filename} (type: {file_type}, size: {len(content)} bytes)")

        try:
            # Images
            if file_type in ("jpg", "jpeg", "png", "gif", "webp", "bmp"):
                assets = await self._process_image(content, filename, context)
            elif file_type in ("tif", "tiff"):
                assets = await self._process_tiff(content, filename, context)
            # PDF - Use hybrid processing (text extraction + vision fallback)
            elif file_type == "pdf":
                if PYMUPDF_AVAILABLE:
                    assets = await self._process_pdf_hybrid(content, filename, context)
                else:
                    assets = await self._process_pdf(content, filename, context)
            # Email
            elif file_type == "msg":
                assets = await self._process_msg(content, filename, context)
            elif file_type == "eml":
                assets = await self._process_eml(content, filename, context)
            # Office documents
            elif file_type == "docx":
                assets = await self._process_docx(content, filename, context)
            elif file_type == "doc":
                # Legacy .doc - try docx parser (may fail for very old formats)
                try:
                    assets = await self._process_docx(content, filename, context)
                except Exception:
                    return ProcessingResult(
                        success=False,
                        error="Legacy .doc format not supported. Please convert to .docx",
                        file_hash=file_hash
                    )
            elif file_type in ("xlsx", "xls"):
                assets = await self._process_xlsx(content, filename, context)
            elif file_type == "pptx":
                assets = await self._process_pptx(content, filename, context)
            elif file_type == "ppt":
                # Legacy .ppt - not supported
                return ProcessingResult(
                    success=False,
                    error="Legacy .ppt format not supported. Please convert to .pptx",
                    file_hash=file_hash
                )
            # Text-based formats
            elif file_type == "txt":
                assets = await self._process_txt(content, filename, context)
            elif file_type == "rtf":
                assets = await self._process_rtf(content, filename, context)
            elif file_type in ("html", "htm"):
                assets = await self._process_html(content, filename, context)
            elif file_type == "csv":
                assets = await self._process_csv(content, filename, context)
            else:
                return ProcessingResult(
                    success=False,
                    error=f"Unsupported file type: {file_type}",
                    file_hash=file_hash
                )

            return ProcessingResult(
                success=True,
                assets=assets,
                file_hash=file_hash
            )

        except Exception as e:
            return ProcessingResult(
                success=False,
                error=str(e),
                file_hash=file_hash
            )

    async def _process_image(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Process an image file, resizing if necessary"""
        logger.info(f"Processing image: {filename}")
        processed_bytes, media_type = self._resize_and_compress_image(content)

        return [ProcessedAsset(
            asset_type="image",
            content=processed_bytes,
            media_type=media_type,
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_pdf(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Convert PDF pages to images - processes one page at a time to save memory"""
        if not PDF2IMAGE_AVAILABLE:
            raise ImportError("pdf2image not available. Install poppler.")

        import gc
        
        logger.info(f"PDF processing started: {filename} ({len(content)} bytes)")

        assets = []
        page_number = 1

        # Process one page at a time to avoid memory exhaustion
        # first_page and last_page are 1-indexed in pdf2image
        while True:
            try:
                # Convert single page at lower DPI to reduce memory usage
                images = convert_from_bytes(
                    content,
                    dpi=100,  # Reduced from 150 to save memory
                    first_page=page_number,
                    last_page=page_number
                )

                if not images:
                    break

                image = images[0]

                # Convert PIL Image to bytes
                buffer = io.BytesIO()
                image.save(buffer, format="JPEG", quality=80, optimize=True)
                image_bytes = buffer.getvalue()

                # Explicitly close and delete image to free memory
                image.close()
                del image
                del images
                buffer.close()

                # Resize if needed
                processed_bytes, media_type = self._resize_and_compress_image(image_bytes)
                del image_bytes

                assets.append(ProcessedAsset(
                    asset_type="image",
                    content=processed_bytes,
                    media_type=media_type,
                    filename=f"{filename}_page_{page_number}.jpg",
                    original_filename=filename,
                    context=context,
                    page_number=page_number
                ))
                
                logger.info(f"PDF page {page_number} converted: {filename}")

                page_number += 1

                # Force garbage collection every few pages
                if page_number % 5 == 0:
                    gc.collect()

            except Exception as e:
                # If we get an error (e.g., page out of range), we're done
                if "page" in str(e).lower() or page_number > 1:
                    break
                raise

        # Final garbage collection
        gc.collect()

        logger.info(f"PDF processing complete: {filename} - {len(assets)} pages converted")

        return assets

    async def _process_pdf_hybrid(
        self,
        content: bytes,
        filename: str,
        context: str = "",
        text_density_threshold: float = 0.001,
        min_text_length: int = 100
    ) -> List[ProcessedAsset]:
        """
        Hybrid PDF processing: extract text from text-based pages, use vision for scanned pages.

        This significantly reduces token usage (55-60% reduction) while maintaining accuracy
        for scanned documents that require OCR via vision.

        Args:
            content: Raw PDF bytes
            filename: Original filename
            context: Additional context
            text_density_threshold: Minimum ratio of text chars to page area (default 0.001)
            min_text_length: Minimum text length to consider page as text-based (default 100)

        Returns:
            List of ProcessedAssets with exact page numbers preserved
        """
        if not PYMUPDF_AVAILABLE:
            logger.warning("PyMuPDF not available, falling back to image-only processing")
            return await self._process_pdf(content, filename, context)

        import gc

        file_size = len(content)
        logger.info(f"Hybrid PDF processing started: {filename} ({file_size} bytes)")

        # Use lower DPI for very large files to reduce memory per page
        image_dpi = 72 if file_size > 5_000_000 else 100

        assets = []
        text_pages = 0
        image_pages = 0

        try:
            doc = fitz.open(stream=content, filetype="pdf")
            total_pages = len(doc)

            if file_size > 5_000_000:
                logger.info(f"Large file detected ({file_size} bytes), using reduced DPI ({image_dpi})")

            for page_num in range(total_pages):
                page = doc[page_num]
                human_page_num = page_num + 1  # 1-indexed for display

                # Extract text from page
                text = page.get_text("text").strip()

                # Calculate text density to detect scanned pages
                rect = page.rect
                page_area = rect.width * rect.height if rect.width > 0 and rect.height > 0 else 1
                text_density = len(text) / page_area

                # Determine if page is text-based or scanned
                is_text_based = len(text) >= min_text_length and text_density >= text_density_threshold

                if is_text_based:
                    # Text-based page: use extracted text (MUCH cheaper in tokens)
                    assets.append(ProcessedAsset(
                        asset_type="text",
                        content=text.encode("utf-8"),
                        media_type="text/plain",
                        filename=f"{filename}_page_{human_page_num}.txt",
                        original_filename=filename,
                        context=context,
                        page_number=human_page_num
                    ))
                    text_pages += 1
                    logger.debug(f"Page {human_page_num}: text-based ({len(text)} chars, density={text_density:.4f})")
                else:
                    # Scanned page: convert to image for vision processing
                    # Process ALL pages - no limits
                    pix = None
                    img_bytes = None
                    processed_bytes = None

                    try:
                        pix = page.get_pixmap(dpi=image_dpi)
                        img_bytes = pix.tobytes("jpeg")

                        # Resize/compress if needed
                        processed_bytes, media_type = self._resize_and_compress_image(img_bytes)

                        assets.append(ProcessedAsset(
                            asset_type="image",
                            content=processed_bytes,
                            media_type=media_type,
                            filename=f"{filename}_page_{human_page_num}.jpg",
                            original_filename=filename,
                            context=context,
                            page_number=human_page_num
                        ))
                        image_pages += 1
                        logger.debug(f"Page {human_page_num}: scanned/image ({len(text)} chars, density={text_density:.4f})")

                    finally:
                        # AGGRESSIVE cleanup after EVERY image page
                        if pix is not None:
                            del pix
                        if img_bytes is not None:
                            del img_bytes
                        # Note: processed_bytes is now in the asset, don't delete
                        gc.collect()

                # Clear page reference
                del page

                # Log progress every 25 pages
                if human_page_num % 25 == 0:
                    logger.info(f"Progress: {human_page_num}/{total_pages} pages ({text_pages} text, {image_pages} image)")

            doc.close()
            del doc

        except Exception as e:
            logger.error(f"Hybrid PDF processing failed for {filename}: {e}")
            # Fall back to image-only processing
            logger.info(f"Falling back to image-only processing for {filename}")
            return await self._process_pdf(content, filename, context)

        gc.collect()

        logger.info(
            f"Hybrid PDF processing complete: {filename} - "
            f"{total_pages} pages ({text_pages} text, {image_pages} image)"
        )

        return assets

    async def process_pdf_chunked(
        self,
        content: bytes,
        filename: str,
        context: str = "",
        chunk_size: int = 30
    ):
        """
        Generator that yields chunks of PDF pages as ProcessedAssets.
        Processes chunk_size pages at a time to avoid memory exhaustion.
        
        Args:
            content: Raw PDF bytes
            filename: Original filename
            context: Additional context
            chunk_size: Number of pages per chunk (default 30)
            
        Yields:
            List[ProcessedAsset] - chunks of processed page images
        """
        if not PDF2IMAGE_AVAILABLE:
            raise ImportError("pdf2image not available. Install poppler.")

        import gc
        
        logger.info(f"PDF chunked processing started: {filename} ({len(content)} bytes), chunk_size={chunk_size}")

        page_number = 1
        chunk_assets = []
        total_pages = 0

        while True:
            try:
                # Convert single page at lower DPI to reduce memory usage
                images = convert_from_bytes(
                    content,
                    dpi=100,
                    first_page=page_number,
                    last_page=page_number
                )

                if not images:
                    break

                image = images[0]

                # Convert PIL Image to bytes
                buffer = io.BytesIO()
                image.save(buffer, format="JPEG", quality=80, optimize=True)
                image_bytes = buffer.getvalue()

                # Explicitly close and delete image to free memory
                image.close()
                del image
                del images
                buffer.close()

                # Resize if needed
                processed_bytes, media_type = self._resize_and_compress_image(image_bytes)
                del image_bytes

                chunk_assets.append(ProcessedAsset(
                    asset_type="image",
                    content=processed_bytes,
                    media_type=media_type,
                    filename=f"{filename}_page_{page_number}.jpg",
                    original_filename=filename,
                    context=context,
                    page_number=page_number
                ))
                
                total_pages += 1
                
                # When chunk is full, yield it and clear memory
                if len(chunk_assets) >= chunk_size:
                    logger.info(f"PDF chunk ready: pages {page_number - chunk_size + 1}-{page_number} of {filename}")
                    yield chunk_assets
                    chunk_assets = []
                    gc.collect()

                page_number += 1

                # Force garbage collection every few pages even within chunk
                if page_number % 10 == 0:
                    gc.collect()

            except Exception as e:
                # If we get an error (e.g., page out of range), we're done
                if "page" in str(e).lower() or page_number > 1:
                    break
                raise

        # Yield any remaining pages
        if chunk_assets:
            logger.info(f"PDF final chunk ready: {len(chunk_assets)} pages of {filename}")
            yield chunk_assets

        # Final garbage collection
        gc.collect()
        logger.info(f"PDF chunked processing complete: {filename} - {total_pages} total pages")

    async def _process_msg(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """
        Process Outlook .msg file recursively.
        Extracts email body and all attachments (including nested emails).
        """
        logger.info(f"Processing MSG email: {filename}")
        assets = []

        # Save to temp file (extract-msg requires file path)
        with tempfile.NamedTemporaryFile(
            suffix=".msg",
            dir=self.temp_dir,
            delete=False
        ) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            msg = extract_msg.Message(tmp_path)

            # Extract email metadata and body
            email_text = self._format_email_body(msg)
            assets.append(ProcessedAsset(
                asset_type="email_body",
                content=email_text.encode("utf-8"),
                media_type="text/plain",
                filename=f"{filename}.txt",
                original_filename=filename,
                context=context
            ))

            # Process attachments
            for attachment in msg.attachments:
                try:
                    att_filename = attachment.longFilename or attachment.shortFilename or "unknown"
                    att_content = attachment.data

                    if att_content is None:
                        continue

                    # Determine attachment context
                    att_context = f"Attachment to email: {msg.subject or filename}"
                    if context:
                        att_context = f"{context} > {att_context}"

                    # Recursive processing for all supported attachments
                    att_ext = Path(att_filename).suffix.lower()
                    if att_ext in self.SUPPORTED_DOC_FORMATS or att_ext in self.SUPPORTED_IMAGE_FORMATS:
                        nested_result = await self.process(
                            att_content,
                            att_filename,
                            att_context
                        )
                        if nested_result.success:
                            for asset in nested_result.assets:
                                asset.parent_filename = filename
                            assets.extend(nested_result.assets)

                except Exception as e:
                    # Log but don't fail entire processing
                    print(f"Failed to process attachment {att_filename}: {e}")
                    continue

            msg.close()

        finally:
            # Clean up temp file
            os.unlink(tmp_path)

        return assets

    async def _process_eml(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Process .eml email file"""
        import email
        from email import policy

        assets = []
        msg = email.message_from_bytes(content, policy=policy.default)

        # Extract email body
        email_text = self._format_email_body_from_eml(msg)
        assets.append(ProcessedAsset(
            asset_type="email_body",
            content=email_text.encode("utf-8"),
            media_type="text/plain",
            filename=f"{filename}.txt",
            original_filename=filename,
            context=context
        ))

        # Process attachments
        for part in msg.walk():
            if part.get_content_maintype() == "multipart":
                continue

            att_filename = part.get_filename()
            if not att_filename:
                continue

            att_content = part.get_payload(decode=True)
            if not att_content:
                continue

            att_context = f"Attachment to email: {msg.get('subject', filename)}"
            if context:
                att_context = f"{context} > {att_context}"

            att_ext = Path(att_filename).suffix.lower()
            if att_ext in self.SUPPORTED_DOC_FORMATS or att_ext in self.SUPPORTED_IMAGE_FORMATS:
                nested_result = await self.process(att_content, att_filename, att_context)
                if nested_result.success:
                    for asset in nested_result.assets:
                        asset.parent_filename = filename
                    assets.extend(nested_result.assets)

        return assets

    async def _process_docx(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Extract text from Word documents (.docx)"""
        logger.info(f"Processing DOCX: {filename}")
        from docx import Document

        doc = Document(io.BytesIO(content))
        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    table_rows.append(row_text)
            if table_rows:
                text_parts.append("\n[Table]\n" + "\n".join(table_rows))

        full_text = "\n\n".join(text_parts)

        return [ProcessedAsset(
            asset_type="text",
            content=full_text.encode("utf-8"),
            media_type="text/plain",
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_xlsx(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Extract text from Excel spreadsheets (.xlsx, .xls)"""
        logger.info(f"Processing XLSX: {filename}")
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        assets = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                # Convert each cell to string, handling None values
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip(" |"):
                    rows.append(row_text)

            if rows:
                sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
                assets.append(ProcessedAsset(
                    asset_type="text",
                    content=sheet_text.encode("utf-8"),
                    media_type="text/plain",
                    filename=f"{filename}_{sheet_name}",
                    original_filename=filename,
                    context=context
                ))

        wb.close()

        # If no sheets had data, return at least one empty asset
        if not assets:
            assets.append(ProcessedAsset(
                asset_type="text",
                content=b"[Empty spreadsheet]",
                media_type="text/plain",
                filename=filename,
                original_filename=filename,
                context=context
            ))

        return assets

    async def _process_pptx(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Extract text from PowerPoint presentations (.pptx)"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        assets = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)

            if texts:
                slide_text = f"[Slide {i}]\n" + "\n".join(texts)
                assets.append(ProcessedAsset(
                    asset_type="text",
                    content=slide_text.encode("utf-8"),
                    media_type="text/plain",
                    filename=f"{filename}_slide_{i}",
                    original_filename=filename,
                    page_number=i,
                    context=context
                ))

        # If no slides had text, return at least one asset
        if not assets:
            assets.append(ProcessedAsset(
                asset_type="text",
                content=b"[Empty presentation]",
                media_type="text/plain",
                filename=filename,
                original_filename=filename,
                context=context
            ))

        return assets

    async def _process_txt(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Process plain text files"""
        # Try UTF-8 first, fallback to latin-1
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        return [ProcessedAsset(
            asset_type="text",
            content=text.encode("utf-8"),
            media_type="text/plain",
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_rtf(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Extract text from RTF files"""
        from striprtf.striprtf import rtf_to_text

        rtf_content = content.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_content)

        return [ProcessedAsset(
            asset_type="text",
            content=text.encode("utf-8"),
            media_type="text/plain",
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_html(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Extract text from HTML files"""
        from bs4 import BeautifulSoup

        # Try to decode with different encodings
        try:
            html_content = content.decode("utf-8")
        except UnicodeDecodeError:
            html_content = content.decode("latin-1")

        soup = BeautifulSoup(html_content, "lxml")

        # Remove script and style elements
        for tag in soup(["script", "style"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)

        return [ProcessedAsset(
            asset_type="text",
            content=text.encode("utf-8"),
            media_type="text/plain",
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_csv(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Process CSV files as formatted text"""
        import csv

        # Try to decode with different encodings
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        # Parse CSV and format as text with pipe separators
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
        formatted = "\n".join(rows)

        return [ProcessedAsset(
            asset_type="text",
            content=formatted.encode("utf-8"),
            media_type="text/plain",
            filename=filename,
            original_filename=filename,
            context=context
        )]

    async def _process_tiff(
        self,
        content: bytes,
        filename: str,
        context: str = ""
    ) -> List[ProcessedAsset]:
        """Process multi-page TIFF images"""
        logger.info(f"Processing TIFF: {filename}")
        import gc

        img = Image.open(io.BytesIO(content))
        assets = []
        page = 0

        while True:
            try:
                img.seek(page)

                # Convert frame to RGB if needed and save as JPEG
                buffer = io.BytesIO()
                frame = img.copy()
                if frame.mode != "RGB":
                    frame = frame.convert("RGB")
                frame.save(buffer, format="JPEG", quality=80, optimize=True)

                processed_bytes, media_type = self._resize_and_compress_image(buffer.getvalue())

                assets.append(ProcessedAsset(
                    asset_type="image",
                    content=processed_bytes,
                    media_type=media_type,
                    filename=f"{filename}_page_{page + 1}.jpg",
                    original_filename=filename,
                    page_number=page + 1,
                    context=context
                ))

                frame.close()
                buffer.close()
                page += 1

                # Garbage collection every few pages
                if page % 5 == 0:
                    gc.collect()

            except EOFError:
                break

        img.close()
        gc.collect()

        return assets

    def _format_email_body(self, msg) -> str:
        """Format extract-msg Message into text"""
        parts = [
            f"From: {msg.sender or 'Unknown'}",
            f"To: {msg.to or 'Unknown'}",
            f"CC: {msg.cc or ''}" if msg.cc else None,
            f"Date: {msg.date or 'Unknown'}",
            f"Subject: {msg.subject or 'No Subject'}",
            "",
            msg.body or ""
        ]
        return "\n".join(p for p in parts if p is not None)

    def _format_email_body_from_eml(self, msg) -> str:
        """Format email.message.Message into text"""
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True).decode("utf-8", errors="replace")
                    break
        else:
            body = msg.get_payload(decode=True).decode("utf-8", errors="replace")

        parts = [
            f"From: {msg.get('from', 'Unknown')}",
            f"To: {msg.get('to', 'Unknown')}",
            f"CC: {msg.get('cc', '')}" if msg.get("cc") else None,
            f"Date: {msg.get('date', 'Unknown')}",
            f"Subject: {msg.get('subject', 'No Subject')}",
            "",
            body
        ]
        return "\n".join(p for p in parts if p is not None)

    def _resize_and_compress_image(
        self,
        image_bytes: bytes,
        max_size_mb: float = None,
        max_dimension: int = None
    ) -> Tuple[bytes, str]:
        """
        Resize and compress image to meet AWS Bedrock constraints.

        Args:
            image_bytes: Raw image bytes
            max_size_mb: Maximum file size in MB
            max_dimension: Maximum width or height

        Returns:
            Tuple of (processed bytes, media type)
        """
        max_size_mb = max_size_mb or self.MAX_IMAGE_SIZE_MB
        max_dimension = max_dimension or self.MAX_IMAGE_DIMENSION
        max_size_bytes = int(max_size_mb * 1024 * 1024)

        img = Image.open(io.BytesIO(image_bytes))

        # Resize if dimensions exceed limit
        if img.width > max_dimension or img.height > max_dimension:
            img.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)

        # Determine format
        original_format = img.format or "JPEG"
        if original_format.upper() in ("JPEG", "JPG"):
            media_type = "image/jpeg"
            save_format = "JPEG"
        elif original_format.upper() == "PNG":
            media_type = "image/png"
            save_format = "PNG"
        elif original_format.upper() == "WEBP":
            media_type = "image/webp"
            save_format = "WEBP"
        else:
            media_type = "image/jpeg"
            save_format = "JPEG"
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")

        # Compress if needed
        buffer = io.BytesIO()

        if save_format in ("JPEG", "WEBP"):
            quality = 95
            while quality > 10:
                buffer.seek(0)
                buffer.truncate()
                img.save(buffer, format=save_format, quality=quality, optimize=True)
                if buffer.tell() <= max_size_bytes:
                    break
                quality -= 5
        else:
            img.save(buffer, format=save_format, optimize=True)

        return buffer.getvalue(), media_type
